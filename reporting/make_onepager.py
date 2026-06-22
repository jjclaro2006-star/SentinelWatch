"""Genera el one-pager de SentinelWatch como un slide 16:9."""
from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

# ── Paleta ────────────────────────────────────────────────────────────────────
DARK    = RGBColor(0x0D, 0x1B, 0x2A)   # azul noche
GREEN   = RGBColor(0x2E, 0xCC, 0x71)   # verde selva
ACCENT  = RGBColor(0xF3, 0x9C, 0x12)   # naranja alerta
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xCC, 0xD6, 0xDD)
MGRAY   = RGBColor(0x1C, 0x2E, 0x3E)   # panel oscuro secundario

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

def rect(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def txbox(slide, l, t, w, h, text, size, color, bold=False,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb

# ── Fondo completo ────────────────────────────────────────────────────────────
rect(slide, 0, 0, W, H, DARK)

# ── Franja izquierda (sidebar) ────────────────────────────────────────────────
SB = Inches(3.4)
rect(slide, 0, 0, SB, H, MGRAY)

# ── Barra verde superior sidebar ──────────────────────────────────────────────
rect(slide, 0, 0, SB, Inches(0.07), GREEN)

# ── Logo / nombre ─────────────────────────────────────────────────────────────
txbox(slide, Inches(0.18), Inches(0.18), SB - Inches(0.18), Inches(0.7),
      "SentinelWatch", 26, GREEN, bold=True)

txbox(slide, Inches(0.18), Inches(0.75), SB - Inches(0.18), Inches(0.4),
      "Detección satelital de minería ilegal en la Amazonía", 8.5, LGRAY, italic=True)

# ── Línea divisora ────────────────────────────────────────────────────────────
rect(slide, Inches(0.18), Inches(1.18), SB - Inches(0.36), Inches(0.025), GREEN)

# ── Métricas clave (sidebar) ──────────────────────────────────────────────────
metrics = [
    ("77.4 %",  "tasa de detección\nvalidación independiente"),
    ("24 / 31", "dragas detectadas\nvs ACCA/Mongabay"),
    ("286 ha",  "área afectada detectada\nRío Nanay, Perú"),
    ("5 capas", "verificación legal\náreas protegidas WDPA"),
]
y = Inches(1.38)
for val, lbl in metrics:
    txbox(slide, Inches(0.18), y, SB - Inches(0.24), Inches(0.38),
          val, 19, ACCENT, bold=True)
    y += Inches(0.38)
    txbox(slide, Inches(0.18), y, SB - Inches(0.24), Inches(0.42),
          lbl, 7.5, LGRAY)
    y += Inches(0.5)

# ── Modelo (sidebar) ──────────────────────────────────────────────────────────
rect(slide, Inches(0.18), Inches(4.95), SB - Inches(0.36), Inches(0.025), ACCENT)
txbox(slide, Inches(0.18), Inches(5.06), SB - Inches(0.24), Inches(0.35),
      "Modelo Gaia v0.5.4", 9, ACCENT, bold=True)
txbox(slide, Inches(0.18), Inches(5.36), SB - Inches(0.24), Inches(1.6),
      "ViT-S/16 (SSL4EO-S12)\n12 bandas Sentinel-2\nFine-tuning supervisado\nThreshold calibrado en campo", 7.5, LGRAY)

# ── Tecnología stack (sidebar base) ───────────────────────────────────────────
txbox(slide, Inches(0.18), Inches(6.82), SB - Inches(0.24), Inches(0.35),
      "Google Earth Engine · PyTorch · FastAPI · GeoJSON", 6.5, LGRAY, italic=True)

# ── Área principal ────────────────────────────────────────────────────────────
MX = SB + Inches(0.35)
MW = W - SB - Inches(0.35)

# Título principal
txbox(slide, MX, Inches(0.22), MW, Inches(0.55),
      "Vigilancia automatizada de la selva en tiempo real", 18, WHITE, bold=True)

rect(slide, MX, Inches(0.72), Inches(2.2), Inches(0.04), GREEN)

# ── Descripción ───────────────────────────────────────────────────────────────
desc = (
    "SentinelWatch es un sistema de detección automatizada de deforestación y minería ilegal "
    "en la Amazonía. Compara imágenes satelitales Sentinel-1 y Sentinel-2 entre períodos de "
    "referencia y análisis, detecta pérdida de vegetación mediante NDVI, y clasifica cada "
    "alerta con un modelo de deep learning (Gaia) entrenado sobre datos reales de campo."
)
txbox(slide, MX, Inches(0.88), MW, Inches(0.95), desc, 8.5, LGRAY, wrap=True)

# ── Tres columnas: Pipeline ───────────────────────────────────────────────────
rect(slide, MX, Inches(1.9), MW, Inches(0.025), MGRAY)
txbox(slide, MX, Inches(1.98), MW, Inches(0.32),
      "PIPELINE DE DETECCIÓN", 7.5, LGRAY, bold=True)

cols = [
    ("01  Adquisición",
     "Google Earth Engine descarga mosaicos S1/S2 del período de referencia y análisis "
     "para la región de interés."),
    ("02  Detección NDVI",
     "Compara índices de vegetación entre períodos. Los píxeles con caída ≥ umbral se "
     "vectorizan como polígonos de alerta."),
    ("03  Clasificación IA",
     "Gaia v0.5.4 (ViT-S/16) infiere si cada chip es minería ilegal, agricultura o "
     "deforestación natural, con probabilidad de confianza."),
    ("04  Verificación legal",
     "Cada alerta se cruza con 5 capas: WDPA, ANP, territorios indígenas, áreas de "
     "amortiguamiento y zonas de reserva forestal."),
]
cw = MW / 4 - Inches(0.12)
cx = MX
for title, body in cols:
    rect(slide, cx, Inches(2.34), cw, Inches(0.04), ACCENT)
    txbox(slide, cx, Inches(2.44), cw, Inches(0.35), title, 8, WHITE, bold=True)
    txbox(slide, cx, Inches(2.8),  cw, Inches(1.15), body, 7.5, LGRAY, wrap=True)
    cx += cw + Inches(0.16)

# ── Sección validación ────────────────────────────────────────────────────────
rect(slide, MX, Inches(4.05), MW, Inches(0.04), GREEN)
txbox(slide, MX, Inches(4.14), MW, Inches(0.32),
      "VALIDACIÓN INDEPENDIENTE — RÍO NANAY, LORETO, PERÚ (JUNIO 2026)", 7.5, GREEN, bold=True)

val_text = (
    "SentinelWatch detectó de forma independiente 24 de 31 dragas documentadas por ACCA/Mongabay "
    "en julio 2024, alcanzando una tasa de detección del 77.4% sin acceso previo a los datos de "
    "referencia. Las 7 dragas no detectadas operan sobre espejo de agua (fuera del alcance de NDVI; "
    "mejora planificada: índice MNDWI). Período de análisis: junio–agosto 2023 vs junio–agosto 2024."
)
txbox(slide, MX, Inches(4.46), MW, Inches(0.85), val_text, 8, LGRAY, wrap=True)

# ── Regiones cubiertas ────────────────────────────────────────────────────────
rect(slide, MX, Inches(5.38), MW, Inches(0.025), MGRAY)
txbox(slide, MX, Inches(5.45), MW * 0.45, Inches(0.28),
      "REGIONES ACTIVAS", 7.5, LGRAY, bold=True)
txbox(slide, MX, Inches(5.73), MW * 0.45, Inches(0.55),
      "Perú (Loreto, Madre de Dios)\nBolivia (Pando, Beni)\nColombia (Amazonas) · Brasil (Acre)",
      8, WHITE, wrap=True)

# ── Output / API ──────────────────────────────────────────────────────────────
txbox(slide, MX + MW * 0.48, Inches(5.45), MW * 0.52, Inches(0.28),
      "SALIDA & API", 7.5, LGRAY, bold=True)
txbox(slide, MX + MW * 0.48, Inches(5.73), MW * 0.52, Inches(0.55),
      "GeoJSON geo-referenciado · FastAPI REST\nMapa web interactivo (index.html)\nAlerta: coordenadas, confianza, estatus legal",
      8, WHITE, wrap=True)

# ── Footer ────────────────────────────────────────────────────────────────────
rect(slide, 0, H - Inches(0.3), W, Inches(0.3), MGRAY)
txbox(slide, Inches(0.2), H - Inches(0.28), W * 0.6, Inches(0.26),
      "SentinelWatch © 2026 — Sistema de monitoreo satelital de la Amazonía", 6.5, LGRAY)
txbox(slide, W * 0.62, H - Inches(0.28), W * 0.36, Inches(0.26),
      "sentinel-1 · sentinel-2 · google earth engine · pytorch · ssl4eo", 6.5, LGRAY,
      align=PP_ALIGN.RIGHT)

out = str(ROOT / "outputs" / "sentinelwatch_onepager.pptx")
prs.save(out)
print(f"Guardado: {out}")
